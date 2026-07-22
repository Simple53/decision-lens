import os
import sys
import json
import argparse
from pathlib import Path

# Try importing dependencies; handled gracefully with uv instructions if missing
try:
    import pandas as pd
except ImportError:
    pd = None

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    Presentation = None


DEFAULT_SAMPLE_DATA = {
    "title": "深入解析 2026 开源桌面 Agent：架构选型与核心矛盾剖析",
    "topic": "开源桌面 Agent 选型与架构深度分析",
    "must_haves": [
        "必须支持离线运行或本地 LLM 接入 (Ollama/LM Studio)",
        "支持工具扩展与 MCP (Model Context Protocol)",
        "具备跨平台支持 (Windows / macOS / Linux)"
    ],
    "sources": [
        {"id": 1, "title": "Anthropic Model Context Protocol Specs", "url": "https://modelcontextprotocol.io"},
        {"id": 2, "title": "GitHub Agent Systems Benchmark 2026", "url": "https://github.com"}
    ],
    "candidates_full": [
        {"name": "Agent A", "type": "桌面独立应用", "passed": True, "reason": "完全符合所有硬性要求", "url": "https://github.com/agent-a"},
        {"name": "Agent B", "type": "VS Code 插件", "passed": False, "reason": "依赖 IDE 宿主，不符合独立桌面运行要求", "url": "https://github.com/agent-b"},
        {"name": "Agent C", "type": "CLI 终端工具", "passed": True, "reason": "支持跨平台与 MCP 扩展", "url": "https://github.com/agent-c"}
    ],
    "core_factors": [
        {
            "num": 1,
            "title": "Model Context Protocol (MCP) 交互架构",
            "desc": "MCP 标准化了 AI Agent 与外部数据源和工具的连接方式。支持 MCP 的 Agent 可以零成本复用社区上百种成熟工具，避免为每个工具单独编写胶水代码。在桌面 Agent 选型中，MCP 原生支持度是决定工具生态丰富度的核心因素。"
        },
        {
            "num": 2,
            "title": "自定义模型接入 (BYOK) 与隐私安全",
            "desc": "桌面 Agent 应当允许用户配置自己的 API Key 或本地模型（如 Ollama / vLLM）。良好的模型接入层不仅能降低使用成本，更重要的是保证敏感数据不流经第三方服务商，实现完全可控的隐私合规。"
        },
        {
            "num": 3,
            "title": "多智能体协作与工作流调度",
            "desc": "复杂任务往往需要规划、执行、代码审查等多个角色分工合作。现代桌面 Agent 逐渐从单 Agent 演化为多 Subagent 架构，主 Agent 负责调度与裁决，子 Agent 在独立上下文中并发执行子任务。"
        }
    ],
    "top_candidates": [
        {
            "name": "Agent A",
            "architecture": "Electron + Rust Core",
            "mcp_support": "原生支持 (Server & Client)",
            "llm_support": "Claude API / OpenAI / Ollama",
            "pros": "UI 现代，沙箱隔离完善，生态活跃",
            "cons": "内存占用较轻量 CLI 偏高"
        },
        {
            "name": "Agent C",
            "architecture": "Python / Asyncio Core",
            "mcp_support": "插件扩展支持",
            "llm_support": "全主流 LLM + Local Models",
            "pros": "极度轻量，定制化脚本扩展能力强",
            "cons": "缺乏图形化可视化界面"
        }
    ],
    "scoring_matrix": [
        {"dimension": "MCP 生态兼容度", "weight": "30%", "agent_a": 5, "agent_c": 4},
        {"dimension": "本地模型与隐私合规", "weight": "30%", "agent_a": 4, "agent_c": 5},
        {"dimension": "易用性与 UI 交互体验", "weight": "20%", "agent_a": 5, "agent_c": 2},
        {"dimension": "资源占用与性能效率", "weight": "20%", "agent_a": 3, "agent_c": 5}
    ],
    "tradeoffs": [
        "如果你最看重交互体验与开箱即用的 UI 界面：Agent A 是更合适的选择，它提供了完善的视窗管理与工具配置，但相比 CLI 占用更多内存空间。",
        "如果你最看重资源效率与高度自定义扩展：Agent C 是最优解，非常适合集成入自动化脚本和低资源环境，但缺乏 GUI。"
    ]
}


def generate_excel(data, output_path):
    """Generate Excel workbook containing full candidates and comparison matrices."""
    if not pd:
        print("Warning: pandas / openpyxl not installed. Skipping Excel generation.")
        return False
    
    with pd.ExcelWriter(output_path, engine='openpyxl') as writer:
        # Sheet 1: Candidate Pool
        df_pool = pd.DataFrame(data.get("candidates_full", []))
        df_pool.to_excel(writer, sheet_name="全量候选池扫描", index=False)
        
        # Sheet 2: Core Matrix
        df_top = pd.DataFrame(data.get("top_candidates", []))
        df_top.to_excel(writer, sheet_name="Top方案深度对比", index=False)
        
        # Sheet 3: Qualitative Scoring
        df_score = pd.DataFrame(data.get("scoring_matrix", []))
        df_score.to_excel(writer, sheet_name="定性评分矩阵(无总分)", index=False)
        
    print(f"✅ Excel report generated at: {output_path}")
    return True


def generate_html(data, output_path):
    """Generate modern, responsive, dark-mode-ready HTML report."""
    html_content = f"""<!DOCTYPE html>
<html lang="zh-CN">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{data.get('title', '决策研究报告')}</title>
    <style>
        :root {{
            --bg-color: #0f172a;
            --card-bg: rgba(30, 41, 59, 0.7);
            --border-color: #334155;
            --primary-accent: #38bdf8;
            --secondary-accent: #818cf8;
            --text-primary: #f8fafc;
            --text-secondary: #94a3b8;
            --tag-pass-bg: #064e3b;
            --tag-pass-fg: #34d399;
            --tag-fail-bg: #7f1d1d;
            --tag-fail-fg: #f87171;
        }}
        body {{
            font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", Roboto, "Helvetica Neue", Arial, sans-serif;
            background-color: var(--bg-color);
            color: var(--text-primary);
            line-height: 1.6;
            margin: 0;
            padding: 2rem;
        }}
        .container {{
            max-width: 1100px;
            margin: 0 auto;
        }}
        header {{
            background: linear-gradient(135deg, rgba(56, 189, 248, 0.1), rgba(129, 140, 248, 0.1));
            border: 1px solid var(--border-color);
            padding: 2.5rem;
            border-radius: 16px;
            margin-bottom: 2rem;
            backdrop-filter: blur(10px);
        }}
        h1 {{
            font-size: 2.2rem;
            margin-top: 0;
            background: linear-gradient(90deg, var(--primary-accent), var(--secondary-accent));
            -webkit-background-clip: text;
            -webkit-text-fill-color: transparent;
        }}
        .section-card {{
            background: var(--card-bg);
            border: 1px solid var(--border-color);
            border-radius: 14px;
            padding: 2rem;
            margin-bottom: 2rem;
            box-shadow: 0 8px 32px 0 rgba(0, 0, 0, 0.36);
        }}
        h2 {{
            color: var(--primary-accent);
            border-bottom: 2px solid var(--border-color);
            padding-bottom: 0.5rem;
            margin-top: 0;
        }}
        .factor-card {{
            background: rgba(15, 23, 42, 0.5);
            border-left: 4px solid var(--secondary-accent);
            padding: 1rem 1.5rem;
            margin-bottom: 1rem;
            border-radius: 0 8px 8px 0;
        }}
        .factor-title {{
            font-weight: bold;
            font-size: 1.1rem;
            color: #e2e8f0;
        }}
        table {{
            width: 100%;
            border-collapse: collapse;
            margin-top: 1rem;
        }}
        th, td {{
            padding: 0.8rem 1rem;
            text-align: left;
            border-bottom: 1px solid var(--border-color);
        }}
        th {{
            background-color: rgba(15, 23, 42, 0.8);
            color: var(--primary-accent);
        }}
        .badge {{
            padding: 0.25rem 0.6rem;
            border-radius: 9999px;
            font-size: 0.85rem;
            font-weight: 600;
        }}
        .badge-pass {{
            background-color: var(--tag-pass-bg);
            color: var(--tag-pass-fg);
        }}
        .badge-fail {{
            background-color: var(--tag-fail-bg);
            color: var(--tag-fail-fg);
        }}
        .tradeoff-item {{
            background: rgba(56, 189, 248, 0.05);
            border: 1px dashed var(--primary-accent);
            padding: 1rem 1.5rem;
            border-radius: 8px;
            margin-bottom: 0.8rem;
        }}
    </style>
</head>
<body>
    <div class="container">
        <header>
            <h1>{data.get('title', '决策分析报告')}</h1>
            <p><strong>研究主题：</strong>{data.get('topic', '')}</p>
        </header>

        <div class="section-card">
            <h2>1. 硬性限制条件 (Must-Have Criteria)</h2>
            <ul>
                {''.join([f"<li>{item}</li>" for item in data.get('must_haves', [])])}
            </ul>
        </div>

        <div class="section-card">
            <h2>2. 10 大核心领域要素详解</h2>
            {''.join([f'''
            <div class="factor-card">
                <div class="factor-title">要素 {f.get('num')}: {f.get('title')}</div>
                <p>{f.get('desc')}</p>
            </div>
            ''' for f in data.get('core_factors', [])])}
        </div>

        <div class="section-card">
            <h2>3. 核心方案对比矩阵</h2>
            <table>
                <thead>
                    <tr>
                        <th>方案名称</th>
                        <th>核心架构</th>
                        <th>MCP 支持</th>
                        <th>LLM 兼容</th>
                        <th>核心优势</th>
                        <th>局限性</th>
                    </tr>
                </thead>
                <tbody>
                    {''.join([f'''
                    <tr>
                        <td><strong>{cand.get('name')}</strong></td>
                        <td>{cand.get('architecture')}</td>
                        <td>{cand.get('mcp_support')}</td>
                        <td>{cand.get('llm_support')}</td>
                        <td>{cand.get('pros')}</td>
                        <td>{cand.get('cons')}</td>
                    </tr>
                    ''' for cand in data.get('top_candidates', [])])}
                </tbody>
            </table>
        </div>

        <div class="section-card">
            <h2>4. 情境权衡分析 (Trade-off Analysis)</h2>
            {''.join([f'<div class="tradeoff-item">💡 {t}</div>' for t in data.get('tradeoffs', [])])}
        </div>
    </div>
</body>
</html>
"""
    with open(output_path, "w", encoding="utf-8") as f:
        f.write(html_content)
    print(f"✅ HTML report generated at: {output_path}")
    return True


def generate_pptx(data, output_path):
    """Generate editable PowerPoint presentation using python-pptx."""
    if not Presentation:
        print("Warning: python-pptx not installed. Skipping PPTX generation.")
        return False
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 widescreen
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # blank layout

    # Slide 1: Title Slide
    slide1 = prs.slides.add_slide(blank_layout)
    # Background shape
    bg = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(15, 23, 42)
    
    txBox = slide1.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11.333), Inches(3))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = data.get("title", "决策研究报告")
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(56, 189, 248)
    
    p2 = tf.add_paragraph()
    p2.text = f"研究主题: {data.get('topic', '')}"
    p2.font.size = Pt(20)
    p2.font.color.rgb = RGBColor(148, 163, 184)

    # Slide 2: Core Factors
    slide2 = prs.slides.add_slide(blank_layout)
    bg2 = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg2.fill.solid()
    bg2.fill.fore_color.rgb = RGBColor(15, 23, 42)

    title_box = slide2.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.7), Inches(1))
    tf2 = title_box.text_frame
    p = tf2.paragraphs[0]
    p.text = "10 大核心领域要素详解"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(56, 189, 248)

    top_pos = Inches(1.8)
    for factor in data.get("core_factors", [])[:3]:
        card = slide2.shapes.add_textbox(Inches(0.8), top_pos, Inches(11.7), Inches(1.5))
        tf_card = card.text_frame
        tf_card.word_wrap = True
        p_title = tf_card.paragraphs[0]
        p_title.text = f"要素 {factor.get('num')}: {factor.get('title')}"
        p_title.font.size = Pt(18)
        p_title.font.bold = True
        p_title.font.color.rgb = RGBColor(129, 140, 248)
        
        p_desc = tf_card.add_paragraph()
        p_desc.text = factor.get("desc", "")
        p_desc.font.size = Pt(14)
        p_desc.font.color.rgb = RGBColor(226, 232, 240)
        
        top_pos += Inches(1.6)

    # Slide 3: Comparison Matrix Table
    slide3 = prs.slides.add_slide(blank_layout)
    bg3 = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg3.fill.solid()
    bg3.fill.fore_color.rgb = RGBColor(15, 23, 42)

    t_box3 = slide3.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.7), Inches(1))
    p = t_box3.text_frame.paragraphs[0]
    p.text = "核心候选方案对比矩阵"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(56, 189, 248)

    cands = data.get("top_candidates", [])
    if cands:
        rows = len(cands) + 1
        cols = 5
        table_shape = slide3.shapes.add_table(rows, cols, Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.5))
        table = table_shape.table
        
        headers = ["方案名称", "核心架构", "MCP 支持", "核心优势", "局限性"]
        for col_idx, text in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = text
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(30, 41, 59)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(14)
                p.font.bold = True
                p.font.color.rgb = RGBColor(56, 189, 248)
                
        for row_idx, cand in enumerate(cands, start=1):
            vals = [cand.get("name"), cand.get("architecture"), cand.get("mcp_support"), cand.get("pros"), cand.get("cons")]
            for col_idx, val in enumerate(vals):
                cell = table.cell(row_idx, col_idx)
                cell.text = str(val)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(15, 23, 42)
                for p in cell.text_frame.paragraphs:
                    p.font.size = Pt(12)
                    p.font.color.rgb = RGBColor(248, 250, 252)

    prs.save(output_path)
    print(f"✅ Editable PowerPoint report generated at: {output_path}")
    return True


def main():
    parser = argparse.ArgumentParser(description="Decision Lens Report Generator (Excel, HTML, PPTX)")
    parser.add_argument("--input", help="Path to input JSON data file", default=None)
    parser.add_argument("--output-dir", help="Directory to save generated reports", default=".")
    parser.add_argument("--name-prefix", help="Prefix for generated files", default="domain_research")
    
    args = parser.parse_args()
    
    if args.input and os.path.exists(args.input):
        with open(args.input, "r", encoding="utf-8") as f:
            data = json.load(f)
    else:
        data = DEFAULT_SAMPLE_DATA
        
    out_dir = Path(args.output_dir)
    out_dir.mkdir(parents=True, exist_ok=True)
    
    excel_path = out_dir / f"{args.name_prefix}_candidates.xlsx"
    html_path = out_dir / f"{args.name_prefix}_report.html"
    pptx_path = out_dir / f"{args.name_prefix}_presentation.pptx"
    
    generate_excel(data, excel_path)
    generate_html(data, html_path)
    generate_pptx(data, pptx_path)


if __name__ == "__main__":
    main()
